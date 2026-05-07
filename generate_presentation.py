"""Generate concise 5-slide PowerPoint for Capstone defense."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

NAVY      = RGBColor(0x0A, 0x14, 0x28)
DARK_BLUE = RGBColor(0x0D, 0x2B, 0x55)
BLUE      = RGBColor(0x14, 0x6E, 0xBE)
BRIGHT    = RGBColor(0x3B, 0x82, 0xF6)
GOLD      = RGBColor(0xF5, 0x9E, 0x0B)
GREEN     = RGBColor(0x10, 0xB9, 0x81)
TEAL      = RGBColor(0x06, 0xB6, 0xD4)
PURPLE    = RGBColor(0x7C, 0x3A, 0xED)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
SLATE     = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT = RGBColor(0x0F, 0x17, 0x2A)
CARD_BG   = RGBColor(0xF1, 0xF5, 0xF9)
LIGHT_BLUE= RGBColor(0xEF, 0xF6, 0xFF)

W = Inches(13.33)
H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def rect(slide, l, t, w, h, color, rounded=False):
    s = slide.shapes.add_shape(5 if rounded else 1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if rounded: s.adjustments[0] = 0.07
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=18, bold=False, color=DARK_TEXT,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def header(slide, title, subtitle, accent=BLUE):
    rect(slide, 0, 0, W, Inches(1.25), accent)
    txt(slide, title, Inches(0.45), Inches(0.1), Inches(11), Inches(0.75),
        size=32, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.45), Inches(0.82), Inches(11), Inches(0.38),
            size=14, color=RGBColor(0xBF,0xDB,0xF7), italic=True)

prs = new_prs()

# ═══════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════
s = blank(prs)
bg(s, NAVY)
rect(s, W-Inches(4.2), 0, Inches(4.2), H, DARK_BLUE)
rect(s, W-Inches(4.2), 0, Inches(0.06), H, BLUE)
rect(s, 0, H-Inches(0.1), W, Inches(0.1), GOLD)
rect(s, 0, 0, W, Inches(0.1), BLUE)

txt(s, "Legal Assistant",
    Inches(0.5), Inches(1.5), Inches(8.3), Inches(1.0),
    size=48, bold=True, color=WHITE)
txt(s, "for SMB Kazakhstan",
    Inches(0.5), Inches(2.45), Inches(8.3), Inches(1.0),
    size=48, bold=True, color=BRIGHT)

rect(s, Inches(0.5), Inches(3.6), Inches(5.5), Inches(0.06), GOLD)

txt(s, "Agentic RAG  |  Multilingual  |  Live Government Data",
    Inches(0.5), Inches(3.75), Inches(8.0), Inches(0.5),
    size=18, color=RGBColor(0x93,0xC5,0xFD), italic=True)

txt(s, "Guldana Kassym-Ashim",
    Inches(0.5), Inches(4.3), Inches(8.0), Inches(0.5),
    size=22, bold=True, color=WHITE)
txt(s, "EPAM  Generative AI for Software Development  May 2026",
    Inches(0.5), Inches(4.85), Inches(8.0), Inches(0.4),
    size=14, color=SLATE)

# Feature badges — 2 rows x 3
badges = [
    ("🔄", "Agentic RAG",      "Retry loop, 3 attempts",    GOLD),
    ("🌐", "MCP Integration",  "Live gov data",              GREEN),
    ("📊", "RAGAS Evaluation", "Quality per response",       BRIGHT),
    ("🌍", "Multilingual",     "Kazakh / English / Russian", TEAL),
    ("👁", "LangFuse",         "Full observability",         PURPLE),
    ("📂", "Custom KB",        "Upload your own docs",       BLUE),
]
bw, bh, bg_gap = Inches(2.78), Inches(0.72), Inches(0.12)
for i, (icon, title, desc, color) in enumerate(badges):
    col, row = i % 3, i // 3
    bl = Inches(0.5) + col * (bw + bg_gap)
    bt = Inches(5.45) + row * (bh + Inches(0.1))
    rect(s, bl, bt, bw, bh, RGBColor(0x12, 0x2A, 0x50), rounded=True)
    rect(s, bl, bt, Inches(0.06), bh, color)
    txt(s, icon,  bl+Inches(0.14), bt+Inches(0.12), Inches(0.4),      Inches(0.45), size=18)
    txt(s, title, bl+Inches(0.58), bt+Inches(0.08), bw-Inches(0.7),   Inches(0.32),
        size=13, bold=True, color=color)
    txt(s, desc,  bl+Inches(0.58), bt+Inches(0.4),  bw-Inches(0.7),   Inches(0.25),
        size=11, color=SLATE)

for i, (num, lbl) in enumerate([("11","Nodes"),("7","AI Agents"),("4","Legal Codes"),("3","Gov Sources")]):
    tp = Inches(1.1) + i * Inches(1.5)
    txt(s, num,  W-Inches(3.8), tp,            Inches(3.5), Inches(0.75),
        size=34, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txt(s, lbl,  W-Inches(3.8), tp+Inches(0.7),Inches(3.5), Inches(0.38),
        size=13, color=SLATE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
#  SLIDE 2 — PROBLEM + COST / ROI
# ═══════════════════════════════════════════════════════════
s = blank(prs)
bg(s, OFF_WHITE)
header(s, "What Problem Does This Solve?", "Legal compliance is expensive and complex for SMBs in Kazakhstan", DARK_BLUE)

# 3 pain cards — compact
pains = [
    ("💸", "High Cost",        GOLD,   "Legal consultations\n$50–200 per hour.\nSMBs can't afford it."),
    ("📚", "Complexity",       BRIGHT, "4 legal codes, 2000+ pages.\nLaws change constantly.\nHard to stay updated."),
    ("🌍", "Language Barrier", GREEN,  "Kazakh, English, Russian.\nNo single accessible source.\nErrors lead to fines."),
]
for i, (icon, title, color, desc) in enumerate(pains):
    l = Inches(0.5) + i * Inches(4.2)
    rect(s, l, Inches(1.35), Inches(3.9), Inches(2.6), NAVY, rounded=True)
    rect(s, l, Inches(1.35), Inches(3.9), Inches(0.07), color)
    txt(s, icon,  l+Inches(0.2),  Inches(1.45), Inches(0.65), Inches(0.55), size=28)
    txt(s, title, l+Inches(0.9),  Inches(1.47), Inches(2.9),  Inches(0.5),
        size=17, bold=True, color=color)
    txt(s, desc,  l+Inches(0.2),  Inches(2.08), Inches(3.55), Inches(1.6),
        size=14, color=RGBColor(0xBF, 0xDB, 0xF7))

# Cost comparison section label
txt(s, "💰  Cost Comparison & ROI", Inches(0.5), Inches(4.12), Inches(7), Inches(0.42),
    size=15, bold=True, color=DARK_TEXT)

ORANGE = RGBColor(0xF9, 0x73, 0x16)
costs = [
    ("👨‍⚖️", "Lawyer (1 hour)",     "$20 – $100",   GOLD,   "Traditional consultation"),
    ("📋",         "PRG / Paragraf (KZ)", "$17 – $45/mo",  ORANGE, "Kazakhstan legal DB subscription"),
    ("🤖",         "Our AI (1 query)",   "~$0.01",        GREEN,  "1,000 queries ≈ $8"),
]
for i, (icon, label, cost, color, sub) in enumerate(costs):
    cl  = Inches(0.5) + i * Inches(4.2)
    ct  = Inches(4.6)
    cw2 = Inches(3.9)
    ch2 = Inches(1.35)
    bg_col = RGBColor(0x05, 0x2E, 0x16) if i == 2 else RGBColor(0x14, 0x1E, 0x3A)
    rect(s, cl, ct, cw2, ch2, bg_col, rounded=True)
    rect(s, cl, ct, Inches(0.06), ch2, color)
    txt(s, icon + "  " + label, cl+Inches(0.15), ct+Inches(0.1),  cw2-Inches(0.25), Inches(0.38),
        size=12, bold=True, color=color)
    txt(s, cost,                cl+Inches(0.15), ct+Inches(0.52), cw2-Inches(0.25), Inches(0.52),
        size=22, bold=True, color=WHITE)
    txt(s, sub,                 cl+Inches(0.15), ct+Inches(1.05), cw2-Inches(0.25), Inches(0.27),
        size=11, color=SLATE, italic=True)

rect(s, Inches(0.5), Inches(6.1), Inches(12.33), Inches(0.72), DARK_BLUE, rounded=True)
txt(s, "ROI: up to 10,000× cheaper than a lawyer  •  Unlimited legal access  •  Answers in seconds, not days",
    Inches(0.7), Inches(6.18), Inches(11.9), Inches(0.55),
    size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
#  SLIDE 3 — TECH STACK
# ═══════════════════════════════════════════════════════════
s = blank(prs)
bg(s, OFF_WHITE)
header(s, "Technology Stack", "Production-grade components — all integrated", BLUE)

stack = [
    ("🧠", "LangGraph",      DARK_BLUE, "Agentic orchestration\n11 nodes, conditional routing\nTypedDict state management"),
    ("🤖", "GPT-4.1-mini",   BLUE,      "LLM for all agents\nGeneration, grading\n& reasoning"),
    ("🔍", "Hybrid Search",  BRIGHT,    "ChromaDB vector search\n+ BM25 keyword search\nRRF Fusion — best of both"),
    ("🌐", "MCP Server",     GREEN,     "Custom FastMCP server\nadilet.zan.kz + kgd.gov.kz\n+ egov.kz live data"),
    ("📊", "RAGAS",          GOLD,      "gpt-4o-mini as judge\nFaithfulness + Relevancy\n~$0.001 per eval"),
    ("👁", "LangFuse",       TEAL,      "Full pipeline tracing\nToken cost tracking\nUser feedback logging"),
    ("🖥", "Gradio UI",      PURPLE,    "Streaming interface\nKazakh / English / Russian\nReal-time agent trace"),
    ("📂", "Knowledge Base",  GREEN,     "Upload your own PDFs & DOCX\n4 Kazakhstan legal codes\nAuto-indexed at runtime"),
]

cw, ch, cg = Inches(3.05), Inches(2.55), Inches(0.15)
for i, (icon, title, color, desc) in enumerate(stack):
    col, row = i%4, i//4
    l = Inches(0.35) + col*(cw+cg)
    t = Inches(1.38) + row*(ch+cg)
    rect(s, l, t, cw, ch, CARD_BG, rounded=True)
    rect(s, l, t, Inches(0.06), ch, color)
    txt(s, icon,  l+Inches(0.15), t+Inches(0.12), Inches(0.5),     Inches(0.5),  size=22)
    txt(s, title, l+Inches(0.68), t+Inches(0.15), cw-Inches(0.83), Inches(0.45),
        size=15, bold=True, color=color)
    txt(s, desc,  l+Inches(0.15), t+Inches(0.7),  cw-Inches(0.3),  ch-Inches(0.85),
        size=12, color=SLATE)


# ═══════════════════════════════════════════════════════════
#  SLIDE 4 — ARCHITECTURE / PIPELINE
# ═══════════════════════════════════════════════════════════
s = blank(prs)
bg(s, OFF_WHITE)
header(s, "Pipeline Architecture", "Agentic RAG with retry loop — 11 nodes, 7 agents", DARK_BLUE)

arch_path = r"c:\Users\Dana\Capstone_EPAM\architecture_diagram.png"
if os.path.exists(arch_path):
    s.shapes.add_picture(arch_path, Inches(0.2), Inches(1.32), Inches(12.93), Inches(5.98))


# ═══════════════════════════════════════════════════════════
#  SLIDE 5 — DEMO
# ═══════════════════════════════════════════════════════════
s = blank(prs)
bg(s, NAVY)
rect(s, 0, 0, W, Inches(0.1), BRIGHT)
txt(s, "Live Demo",
    Inches(0.5), Inches(0.12), Inches(12.33), Inches(0.72),
    size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

ui_path = r"c:\Users\Dana\Capstone_EPAM\ui_screenshot.png"
if os.path.exists(ui_path):
    s.shapes.add_picture(ui_path, Inches(0.2), Inches(0.95), Inches(12.93), Inches(6.42))
else:
    rect(s, Inches(0.2), Inches(0.95), Inches(12.93), Inches(6.42), DARK_BLUE, rounded=True)
    txt(s, "[ Insert UI Screenshot ]", Inches(4), Inches(3.7), Inches(5), Inches(0.7),
        size=24, color=SLATE, align=PP_ALIGN.CENTER, italic=True)


# ── Save ────────────────────────────────────────────────────
out = r"c:\Users\Dana\Capstone_EPAM\Capstone_Presentation.pptx"
prs.save(out)
print("Saved: " + out)
print("Slides: " + str(len(prs.slides)))
